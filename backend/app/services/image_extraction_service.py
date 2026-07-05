import os
import fitz
import cv2

from PIL import Image
from app.services.image.image_quality_service import (
    is_good_embedded_image
)
from app.services.diagram_detector import detect_diagrams


# ---------------------------------------------------------
# Folder helper
# ---------------------------------------------------------

def create_image_folder(document_id):

    folder = f"uploads/images/{document_id}"

    os.makedirs(folder, exist_ok=True)

    return folder


# ---------------------------------------------------------
# PDF Extraction
# ---------------------------------------------------------

def extract_pdf_images(pdf_path, document_id):

    pdf = fitz.open(pdf_path)

    output_folder = create_image_folder(document_id)

    extracted = []
    embedded_count = 0
    filtered_count = 0
    diagram_count = 0

    for page_no in range(len(pdf)):

        page = pdf.load_page(page_no)

        # =====================================================
        # Render page at higher resolution
        # =====================================================

        pix = page.get_pixmap(
    matrix=fitz.Matrix(2, 2)
)

        page_path = os.path.join(
            output_folder,
            f"page_{page_no}.png"
        )

        pix.save(page_path)

        # =====================================================
        # Extract embedded images
        # =====================================================

        images = page.get_images(full=True)

        for img in images:

            try:

                xref = img[0]

                base = pdf.extract_image(xref)

                image_bytes = base["image"]

                ext = base["ext"]

                save_path = os.path.join(
                    output_folder,
                    f"embedded_{page_no}_{xref}.{ext}"
                )

                with open(save_path, "wb") as f:
                    f.write(image_bytes)

                if not is_good_embedded_image(save_path):
                    filtered_count += 1

                    os.remove(save_path)

                    continue
                embedded_count += 1

                extracted.append(

    {

        "path": save_path,

        "page_no": page_no + 1,

        "source": "embedded",

        "bbox": None

    }

)
                

            except Exception:
                continue

        # =====================================================
        # Detect diagrams from rendered page
        # (Always run, even if embedded images exist)
        # =====================================================

        diagrams = detect_diagrams(page_path)

        for i, diagram in enumerate(diagrams):

            save_path = os.path.join(

                output_folder,

                f"diagram_{page_no}_{i}.png"

            )

            cv2.imwrite(

                save_path,

                diagram["image"]

            )
            diagram_count += 1

            extracted.append({

                "path": save_path,

                "page_no": page_no + 1,

                "source": "detected",

                "bbox": diagram.get("bbox")

            })
    print("\n========== Image Extraction ==========")
    print(f"Embedded Images : {embedded_count}")
    print(f"Filtered Images : {filtered_count}")
    print(f"Detected Figures: {diagram_count}")
    print(f"Total Images    : {len(extracted)}")
    print("======================================\n")

    return extracted


# ---------------------------------------------------------
# DOCX
# ---------------------------------------------------------

def extract_docx_images(file_path, document_id):

    from docx import Document

    folder = create_image_folder(document_id)

    extracted = []

    doc = Document(file_path)

    index = 0

    for rel in doc.part._rels.values():

        if "image" in rel.target_ref:

            index += 1
            image_data = rel.target_part.blob

            path = os.path.join(

                folder,

                f"docx_{index}.png"

            )

            with open(path, "wb") as f:

                f.write(image_data)

            extracted.append({

                "path": path,

                "page_no": 1,

                "source": "embedded",

                "bbox": None

            })

    return extracted


# ---------------------------------------------------------
# PPTX
# ---------------------------------------------------------

def extract_pptx_images(file_path, document_id):

    from pptx import Presentation

    folder = create_image_folder(document_id)

    prs = Presentation(file_path)

    extracted = []

    index = 0

    for slide_no, slide in enumerate(prs.slides):

        for shape in slide.shapes:

            if hasattr(shape, "image"):

                try:

                    image = shape.image

                    index += 1

                    path = os.path.join(

                        folder,

                        f"slide_{slide_no}_{index}.{image.ext}"

                    )

                    with open(path, "wb") as f:

                        f.write(image.blob)

                    extracted.append({

                        "path": path,

                        "page_no": slide_no + 1,

                        "source": "embedded",

                        "bbox": None

                    })

                except Exception:
                    pass

    return extracted


# ---------------------------------------------------------
# Single Image Upload
# ---------------------------------------------------------

def extract_image_file(file_path):

    return [

        {

            "path": file_path,

            "page_no": 1,

            "source": "uploaded",

            "bbox": None

        }

    ]


# ---------------------------------------------------------
# Generic
# ---------------------------------------------------------

def extract_images(file_path, document_id):

    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pdf":

        return extract_pdf_images(file_path, document_id)

    elif extension == ".docx":

        return extract_docx_images(file_path, document_id)

    elif extension == ".pptx":

        return extract_pptx_images(file_path, document_id)

    elif extension in [".png", ".jpg", ".jpeg"]:

        return extract_image_file(file_path)

    return []
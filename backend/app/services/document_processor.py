import os
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import fitz  # PyMuPDF


import fitz
import re


def split_into_chunks(text, chunk_size=800):

    text = re.sub(r"\n{2,}", "\n\n", text)

    paragraphs = text.split("\n\n")

    chunks = []

    current = ""

    for para in paragraphs:

        para = para.strip()

        if not para:
            continue

        if len(current) + len(para) < chunk_size:

            current += para + "\n\n"

        else:

            if current.strip():
                chunks.append(current.strip())

            current = para + "\n\n"

    if current.strip():
        chunks.append(current.strip())

    return chunks


def extract_pdf_text(file_path):

    doc = fitz.open(file_path)

    pages = []

    for page_num in range(len(doc)):

        page = doc.load_page(page_num)

        text = page.get_text("text")

        chunks = split_into_chunks(text)

        if len(chunks) == 0:

            chunks = [""]

        for chunk in chunks:

            pages.append({

                "page_no": page_num + 1,

                "text": chunk

            })

    return pages
def extract_docx_text(file_path):
    doc = Document(file_path)

    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return [{
        "page_no": 1,
        "text": text
    }]


def extract_pptx_text(file_path):
    prs = Presentation(file_path)

    slides = []

    for slide_no, slide in enumerate(prs.slides):

        text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        slides.append({
            "page_no": slide_no + 1,
            "text": text
        })

    return slides


def extract_image_text(file_path):
    pytesseract.pytesseract.tesseract_cmd = (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    )

    image = Image.open(file_path)

    text = pytesseract.image_to_string(image)

    return [{
        "page_no": 1,
        "text": text
    }]


def extract_txt_text(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    return [{
        "page_no": 1,
        "text": text
    }]


def extract_text(file_path):

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_text(file_path)

    elif ext == ".txt":
        return extract_txt_text(file_path)

    elif ext == ".docx":
        return extract_docx_text(file_path)

    elif ext == ".pptx":
        return extract_pptx_text(file_path)

    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_image_text(file_path)

    else:
        raise Exception("Unsupported file type")